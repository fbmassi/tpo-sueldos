"""Presentacion de defensa — estructura estricta del usuario, estetica B&N alto contraste."""
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
EDA = ROOT/"data/processed/eda"; MOD = ROOT/"data/processed/modelos"; PRE = ROOT/"data/processed/presentacion"
OUT = ROOT/"presentacion/defensa_oral.pptx"

BLACK=RGBColor(0x0A,0x0A,0x0A); WHITE=RGBColor(0xFF,0xFF,0xFF); ACC=RGBColor(0xE6,0xFF,0x00)
GRAY=RGBColor(0x8C,0x8C,0x8C); LGRAY=RGBColor(0xED,0xED,0xED); DCARD=RGBColor(0x1A,0x1A,0x1A)
FH="Arial"; FB="Calibri"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]

def slide(bg):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s

def box(s,l,t,w,h):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=tf.margin_right=Inches(0.05); tf.margin_top=tf.margin_bottom=Inches(0.02)
    return tb,tf

def P(tf,txt,size,color,bold=False,italic=False,align=PP_ALIGN.LEFT,sa=6,bullet=False,first=False,font=FB):
    p=tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(sa); p.space_before=Pt(0)
    if bullet: txt="—  "+txt
    r=p.add_run(); r.text=txt; r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color; r.font.name=font
    return p

def rect(s,l,t,w,h,color,rounded=False,line=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                          Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=color
    if line: sh.line.color.rgb=line; sh.line.width=Pt(1.5)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; return sh

def kicker(s,num,sec,dark):
    _,tf=box(s,0.6,0.42,11,0.5)
    p=tf.paragraphs[0]; p.space_after=Pt(0)
    r=p.add_run(); r.text=f"{num}  ·  "; r.font.size=Pt(13); r.font.bold=True; r.font.name=FH; r.font.color.rgb=ACC if dark else BLACK
    r2=p.add_run(); r2.text=sec.upper(); r2.font.size=Pt(13); r2.font.bold=True; r2.font.name=FH; r2.font.color.rgb=GRAY

def title(s,txt,dark,size=34):
    _,tf=box(s,0.6,0.95,12.1,1.0)
    P(tf,txt,size,WHITE if dark else BLACK,bold=True,first=True,font=FH,sa=0)

def img_fit(s,path,l,t,w,h):
    iw,ih=Image.open(path).size; ar=iw/ih; bar=w/h
    nw,nh=(w,w/ar) if ar>bar else (h*ar,h)
    s.shapes.add_picture(str(path),Inches(l+(w-nw)/2),Inches(t+(h-nh)/2),Inches(nw),Inches(nh))

def stat(s,l,t,w,num,lab,dark,numc=None):
    _,tf=box(s,l,t,w,1.5)
    P(tf,num,46,numc or (ACC if dark else BLACK),bold=True,align=PP_ALIGN.CENTER,first=True,sa=0,font=FH)
    if lab: P(tf,lab,12,GRAY,align=PP_ALIGN.CENTER)

# ============ 1 PORTADA
s=slide(BLACK)
rect(s,0.6,2.0,1.3,0.18,ACC)
_,tf=box(s,0.55,2.35,12.2,3.0)
P(tf,"MERCADO LABORAL",58,WHITE,bold=True,first=True,font=FH,sa=2)
p=tf.add_paragraph(); p.space_after=Pt(6)
r=p.add_run(); r.text="TECH "; r.font.size=Pt(58); r.font.bold=True; r.font.name=FH; r.font.color.rgb=WHITE
r=p.add_run(); r.text="ARGENTINO"; r.font.size=Pt(58); r.font.bold=True; r.font.name=FH; r.font.color.rgb=ACC
P(tf,"Análisis y predicción de salarios del sector IT",20,GRAY,sa=0)
_,tf=box(s,0.6,6.4,12.1,0.7)
P(tf,"TPO Minería de Datos   ·   Sysarmy 2022–2025   ·   31.088 profesionales",13,GRAY,first=True,font=FH)

# ============ 2 EQUIPO Y ESTRATEGIA
s=slide(BLACK); kicker(s,"00","Equipo y estrategia",True); title(s,"Quiénes somos y cómo trabajamos",True)
roles=[("Datos / ETL","Descarga, integración de 9 fuentes y limpieza","Nombre: __________"),
       ("Feature eng. / Modelado","Variables, Random Forest, K-Means y evaluación","Nombre: __________"),
       ("EDA / Visualización","Análisis exploratorio y storytelling","Nombre: __________"),
       ("App / Documentación","Streamlit, informe y presentación","Nombre: __________")]
x=0.6
for t,d,n in roles:
    rect(s,x,2.05,2.95,3.0,DCARD,rounded=True)
    _,tf=box(s,x+0.2,2.25,2.6,2.7)
    P(tf,t,15,ACC,bold=True,first=True,font=FH,sa=6); P(tf,d,12.5,WHITE,sa=10); P(tf,n,11.5,GRAY)
    x+=3.07
_,tf=box(s,0.6,5.35,12.1,1.7)
P(tf,"Estrategia (metodología CRISP-DM):",15,ACC,bold=True,first=True,font=FH,sa=6)
P(tf,"Comprensión del negocio → Datos → Preparación (limpieza + feature engineering) → "
     "Modelado → Evaluación → Despliegue. Trabajo iterativo y versionado en Git, con ramas por feature.",
  13.5,WHITE)

# ============ 3 DOMINIO
s=slide(WHITE); kicker(s,"01","Descripción del dominio y problemática",False)
title(s,"Mercado laboral tech argentino",False)
_,tf=box(s,0.6,1.95,6.5,4.9)
P(tf,"El contexto",16,BLACK,bold=True,first=True,font=FH,sa=8)
for t in ["Sector dinámico, alta rotación y fuerte demanda internacional.",
          "Salarios atados a una macro volátil: +900% de inflación en 3 años y un dólar que se cuadruplicó.",
          "Falta de información integrada sobre salarios: cada encuesta usa otra moneda y otro nivel de precios."]:
    P(tf,t,14.5,BLACK,bullet=True,sa=11)
rect(s,7.4,1.95,5.3,4.7,BLACK,rounded=True)
_,tf=box(s,7.7,2.2,4.7,4.3)
P(tf,"HIPÓTESIS CENTRAL",14,ACC,bold=True,first=True,font=FH,sa=10)
P(tf,"El sueldo real de un profesional tech está determinado por su perfil (rol, experiencia) "
     "y su contexto (provincia, modalidad, tamaño de empresa), una vez expresado en moneda "
     "constante para que las ediciones sean comparables entre sí.",15,WHITE,italic=True)

# ============ 4 HIPOTESIS DESAGREGADAS
s=slide(BLACK); kicker(s,"01","Hipótesis a testear",True); title(s,"Seis hipótesis desagregadas",True)
hs=[("H1","Seniority/experiencia es el factor más predictivo"),
    ("H2","La experiencia tiene rendimientos decrecientes"),
    ("H3","Hay tecnologías con premium salarial"),
    ("H4","CABA paga más, pero menos de lo que se cree"),
    ("H5","Cobrar en dólares da más estabilidad"),
    ("H6","El ITCRM mueve los sueldos en dólares")]
y=2.0
for i,(h,t) in enumerate(hs):
    col=0.6 if i<3 else 6.9; yy=2.15+(i%3)*1.5 if i<3 else 2.15+(i-3)*1.5
    rect(s,col,yy,0.95,0.92,ACC)
    _,tf=box(s,col,yy+0.18,0.95,0.7); P(tf,h,22,BLACK,bold=True,align=PP_ALIGN.CENTER,first=True,font=FH,sa=0)
    _,tf=box(s,col+1.15,yy+0.22,4.6,0.7); P(tf,t,14,WHITE,first=True,sa=0)

# ============ 5 PROPUESTA Y VALOR
s=slide(WHITE); kicker(s,"02","Propuesta y valor para el negocio",False)
title(s,"Qué construimos y para qué sirve",False)
_,tf=box(s,0.6,1.95,6.4,4.9)
P(tf,"Pipeline de datos",16,BLACK,bold=True,first=True,font=FH,sa=8)
for t in ["Unifica 6 ediciones de Sysarmy con 8 series macro.",
          "Produce un dataset único, limpio y comparable en el tiempo.",
          "Alimenta los modelos y una app web interactiva."]:
    P(tf,t,14.5,BLACK,bullet=True,sa=10)
P(tf,"Valor para la gerencia",16,BLACK,bold=True,sa=8)
for t in ["Benchmarking salarial objetivo para ofertas y retención.",
          "Detección de brechas y de los factores que mueven el sueldo.",
          "Lectura del impacto macro sobre el costo del talento."]:
    P(tf,t,14.5,BLACK,bullet=True,sa=10)
rect(s,7.4,1.95,5.3,4.7,ACC,rounded=True)
_,tf=box(s,7.7,2.25,4.7,4.1)
P(tf,"EN NÚMEROS",13,BLACK,bold=True,first=True,font=FH,sa=14)
for n,l in [("31.088","profesionales analizados"),("6","ediciones unificadas"),
            ("9","fuentes de datos"),("2","apps (ES/ARS · EN/USD)")]:
    pp=tf.add_paragraph(); pp.space_after=Pt(10)
    r=pp.add_run(); r.text=n+"   "; r.font.size=Pt(26); r.font.bold=True; r.font.name=FH; r.font.color.rgb=BLACK
    r=pp.add_run(); r.text=l; r.font.size=Pt(13); r.font.color.rgb=DCARD
P(tf,"",10,BLACK)

# ============ 6 FUENTES
s=slide(WHITE); kicker(s,"02","Fuentes utilizadas",False); title(s,"Nueve fuentes integradas",False)
fuentes=[("Sysarmy (6 ed.)","Salarios y perfiles — 31.088 resp."),("IPC INDEC","Inflación AR (deflactor pesos)"),
         ("Dólar MEP / Blue","Tipo de cambio histórico"),("US CPI (FRED)","Inflación EE.UU. (deflactor USD)"),
         ("CBT INDEC","Canasta básica (poder adquis.)"),("RIPTE","Salario formal promedio"),
         ("Big Mac Index","Paridad de poder de compra"),("ITCRM (BCRA)","Tipo de cambio real multilateral"),
         ("Stack Overflow","Benchmark internacional")]
i=0
for r_ in range(3):
    for c_ in range(3):
        x=0.6+c_*4.1; y=2.0+r_*1.55; t,d=fuentes[i]; i+=1
        rect(s,x,y,3.9,1.35,LGRAY,rounded=True)
        _,tf=box(s,x+0.2,y+0.16,3.55,1.1)
        P(tf,t,14,BLACK,bold=True,first=True,font=FH,sa=3); P(tf,d,11.5,GRAY)

# ============ 7 PIPELINE
s=slide(BLACK); kicker(s,"03","Arquitectura de la solución",True); title(s,"Pipeline — diagrama de flujo",True)
steps=[("1","DESCARGA","9 fuentes → data/raw"),("2","LIMPIEZA +\nUNIFICACIÓN","6 ediciones armonizadas"),
       ("3","AJUSTE REAL","IPC + US CPI · base may-2026"),("4","DATASET\nFINAL","31.088 × 17 + contexto")]
x=0.7
for n,t,d in steps:
    rect(s,x,2.3,2.7,2.4,DCARD,rounded=True)
    _,tf=box(s,x+0.2,2.5,2.35,2.1)
    P(tf,n,30,ACC,bold=True,first=True,font=FH,sa=4); P(tf,t,14.5,WHITE,bold=True,font=FH,sa=5); P(tf,d,11.5,GRAY)
    if x<10:
        a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x+2.73),Inches(3.2),Inches(0.32),Inches(0.5))
        a.fill.solid(); a.fill.fore_color.rgb=ACC; a.line.fill.background(); a.shadow.inherit=False
    x+=3.07
_,tf=box(s,0.7,5.2,11.9,1.6)
P(tf,"Salidas: dataset_final + contexto_macroeconomico (FK fecha_edicion) · 12 láminas de EDA · "
     "modelos (RF + K-Means) · app Streamlit (ES/ARS y EN/USD).",13.5,GRAY,first=True)

# ============ 8 LIMPIEZA (foco)
s=slide(WHITE); kicker(s,"03","Foco: limpieza de datos",False); title(s,"De 6 encuestas heterogéneas a un dataset",False)
_,tf=box(s,0.6,1.95,7.4,4.9)
for t,d in [("Armonización de esquemas","cada edición traía 43–56 columnas con nombres distintos; mapeo por tokens al mismo esquema."),
            ("Parseo y tipado","salario de texto a número; fechas; normalización de texto."),
            ("Valores faltantes","eliminación en columnas críticas (provincia, salario); imputación por mediana/moda en el resto."),
            ("Outliers","se conservan marcados; sólo se descartan errores de carga evidentes (>50× la mediana).")]:
    P(tf,t,15,BLACK,bold=True,first=(t==tf.paragraphs[0].text),font=FH,sa=2)
    P(tf,d,13,BLACK,sa=12)
rect(s,8.3,1.95,4.4,4.7,BLACK,rounded=True)
_,tf=box(s,8.6,2.4,3.9,4.0)
stat(s,8.4,2.5,4.2,"32.309",None,True); _,tf2=box(s,8.4,3.45,4.2,0.5)
P(tf2,"filas crudas",12,GRAY,align=PP_ALIGN.CENTER,first=True)
a=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,Inches(10.2),Inches(3.95),Inches(0.4),Inches(0.55))
a.fill.solid(); a.fill.fore_color.rgb=ACC; a.line.fill.background(); a.shadow.inherit=False
stat(s,8.4,4.55,4.2,"31.088",None,True); _,tf2=box(s,8.4,5.5,4.2,0.6)
P(tf2,"filas finales  ·  sólo 3,8% descartado",12,GRAY,align=PP_ALIGN.CENTER,first=True)

# ============ 8b LA LIMPIEZA EN NUMEROS
s=slide(WHITE); kicker(s,"03","Foco: limpieza de datos",False); title(s,"La limpieza en números, edición por edición",False)
img_fit(s,PRE/"limpieza_filas_bn.png",0.4,1.85,8.6,4.9)
rect(s,9.2,1.95,3.6,4.6,BLACK,rounded=True)
_,tf=box(s,9.45,2.2,3.15,4.2)
P(tf,"QUÉ MUESTRA",13,ACC,bold=True,first=True,font=FH,sa=8)
for t in ["Gris: lo que llegó crudo","Negro: lo que quedó tras limpiar",
          "La brecha es mínima en todas las ediciones",
          "Se descarta sólo lo crítico irrecuperable (sin salario o provincia) y errores de carga"]:
    P(tf,t,12,WHITE,bullet=True,sa=9)
P(tf,"Limpiar ≠ recortar: es conservar datos confiables.",12,ACC,italic=True)

# ============ 9 FEATURE ENGINEERING (agregado)
s=slide(BLACK); kicker(s,"03","Foco: feature engineering",True); title(s,"Construcción de variables — y su justificación",True)
fe=[("Tecnologías: FUERA del modelo","su premium era proxy del contexto (backend/dólar) y como feature generaban ruido; se analizan en el EDA."),
    ("Agrupar roles (698 → top 15 + «Otro»)","reduce cardinalidad y evita dummies ruidosas / overfitting."),
    ("Agrupar provincias (<100 → «Otra»)","misma lógica: categorías raras no aportan señal estable."),
    ("Target real (deflación IPC / US CPI)","ingeniería del objetivo: hace los salarios comparables en el tiempo."),
    ("Selección: se quita seniority","redundante con experiencia (multicolinealidad) → mismo R² con menos ruido."),
    ("Anti-leakage","se excluyen variables derivadas del salario (USD, canastas).")]
y=2.0
for i,(t,d) in enumerate(fe):
    col=0.6 if i<3 else 6.85; yy=2.0+(i%3)*1.55
    rect(s,col,yy,6.0,1.4,DCARD,rounded=True)
    _,tf=box(s,col+0.2,yy+0.16,5.65,1.15)
    P(tf,t,13.5,ACC,bold=True,first=True,font=FH,sa=3); P(tf,d,12,WHITE)

# ============ 9b ROLES EN LAS ENCUESTAS (ejemplo de agrupación)
s=slide(WHITE); kicker(s,"03","Foco: feature engineering",False); title(s,"El caso de los roles: 698 → 16 categorías",False)
img_fit(s,PRE/"roles_pareto_bn.png",0.4,1.85,8.6,4.9)
rect(s,9.2,1.95,3.6,4.6,BLACK,rounded=True)
_,tf=box(s,9.45,2.2,3.15,4.2)
P(tf,"POR QUÉ AGRUPAR",13,ACC,bold=True,first=True,font=FH,sa=8)
for t in ["698 puestos distintos escritos de mil maneras",
          "El top-15 cubre el 90% de los casos",
          "683 roles suman apenas el 10% restante",
          "Menos categorías raras → menos columnas dummy → menos overfitting"]:
    P(tf,t,12,WHITE,bullet=True,sa=9)
P(tf,"La curva ámbar es la cobertura acumulada.",11.5,ACC,italic=True)

# ============ 9c IMPACTO DE SACAR LAS TECNOLOGIAS
s=slide(WHITE); kicker(s,"03","Foco: feature engineering",False); title(s,"La prueba: el modelo con y sin tecnologías",False)
img_fit(s,PRE/"fe_impacto_bn.png",0.4,1.85,8.6,4.9)
rect(s,9.2,1.95,3.6,4.6,BLACK,rounded=True)
_,tf=box(s,9.45,2.2,3.15,4.2)
P(tf,"LA EVIDENCIA",13,ACC,bold=True,first=True,font=FH,sa=8)
for t in ["Entrenamos las dos versiones y comparamos",
          "R² casi igual: 0,297 → 0,272",
          "Menos sobreajuste (gap train−test más chico)",
          "Y sin disparates: «Go» en un helpdesk ya no infla la estimación"]:
    P(tf,t,12,WHITE,bullet=True,sa=9)
P(tf,"Decisión basada en evidencia, no en intuición.",11.5,ACC,italic=True)

# ============ 10 EDA distribucion
s=slide(WHITE); kicker(s,"04","EDA",False); title(s,"Distribución de salarios",False)
img_fit(s,PRE/"distribucion_bn.png",0.5,1.9,9.3,4.6)
stat(s,9.9,2.1,3.0,"$3,19M","mediana (ARS)",False)
stat(s,9.9,3.7,3.0,"USD 2.264","mediana (USD)",False)
stat(s,9.9,5.3,3.0,"31.088","profesionales",False)

# ============ 11 EDA seniority/geo
s=slide(WHITE); kicker(s,"04","EDA",False); title(s,"Seniority, modalidad y provincias",False)
img_fit(s,EDA/"eda_02_seniority.png",0.5,1.85,6.2,4.5)
img_fit(s,EDA/"eda_05_geografia.png",6.8,1.85,6.1,4.5)
_,tf=box(s,0.6,6.45,12.2,0.8)
P(tf,"Seniority casi sin solapamiento (jr $1,8M → sr $4,3M)  ·  CABA +24% vs interior  ·  remoto +43% en USD.",
  13.5,BLACK,first=True,align=PP_ALIGN.CENTER,italic=True)

# ============ 12 EDA tecnologias
s=slide(WHITE); kicker(s,"04","EDA",False); title(s,"Tecnologías más usadas",False)
img_fit(s,EDA/"eda_04_tecnologias.png",0.5,1.85,8.3,4.8)
rect(s,9.1,2.05,3.7,4.4,BLACK,rounded=True)
_,tf=box(s,9.35,2.3,3.2,4.0)
P(tf,"EL PREMIUM",14,ACC,bold=True,first=True,font=FH,sa=8)
P(tf,"No lo da la cantidad de tecnologías, sino cuáles:",13,WHITE,sa=10)
for t in ["Go +45%  ·  Rust +48%  ·  Scala +42%","Python +13%","Frontend web (CSS/HTML) bajo la media"]:
    P(tf,t,13,WHITE,bullet=True,sa=9)

# ============ 13 DISPERSION + EJES
s=slide(BLACK); kicker(s,"04","EDA · dispersión y variables",True); title(s,"Cómo se eligieron los ejes (PCA)",True)
img_fit(s,PRE/"clusters_bn.png",0.5,1.9,8.4,4.5)
rect(s,9.15,2.0,3.65,4.5,DCARD,rounded=True)
_,tf=box(s,9.4,2.25,3.15,4.1)
P(tf,"PROYECCIÓN 2D",13,ACC,bold=True,first=True,font=FH,sa=8)
P(tf,"El perfil tiene muchas dimensiones; PCA las comprime a 2 ejes:",12.5,WHITE,sa=8)
P(tf,"Comp. 1 (48%) — años de carrera (junior → senior)",12.5,WHITE,bullet=True,sa=8)
P(tf,"Comp. 2 (15%) — rotación / antigüedad",12.5,WHITE,bullet=True,sa=10)
P(tf,"Puntos cercanos = perfiles parecidos.",12.5,ACC,italic=True)

# ============ 14 VARIABLE TARGET
s=slide(WHITE); kicker(s,"05","Modelos de minería de datos",False); title(s,"Variable target",False)
_,tf=box(s,0.6,2.0,6.4,4.6)
P(tf,"Qué predecimos",16,BLACK,bold=True,first=True,font=FH,sa=8)
P(tf,"El salario bruto REAL, no el nominal: deflactado a pesos (por IPC) y a dólares (por US CPI) "
     "de mayo 2026. Así un perfil es comparable entre ediciones.",14.5,BLACK,sa=12)
P(tf,"Por qué real y no nominal",16,BLACK,bold=True,font=FH,sa=8)
P(tf,"El número nominal creció 10× sólo por inflación: entrenar sobre nominal mezclaría épocas y "
     "aprendería el paso del tiempo, no el perfil.",14.5,BLACK)
rect(s,7.4,2.0,5.3,4.5,BLACK,rounded=True)
_,tf=box(s,7.75,2.5,4.6,3.6)
P(tf,"DOS UNIDADES",13,ACC,bold=True,first=True,font=FH,sa=12)
P(tf,"salario_real_ars",17,WHITE,bold=True,font=FH,sa=2); P(tf,"pesos constantes may-2026",12,GRAY,sa=14)
P(tf,"salario_real_usd",17,WHITE,bold=True,font=FH,sa=2); P(tf,"dólares constantes may-2026",12,GRAY,sa=10)
P(tf,"Difieren en una constante → mismo R².",12.5,ACC,italic=True)

# ============ 15 COMPARACION DE MODELOS
s=slide(WHITE); kicker(s,"05","Comparación de modelos",False); title(s,"Cinco técnicas comparadas",False)
img_fit(s,PRE/"modelos_bn.png",0.5,1.85,8.0,4.6)
rect(s,8.7,2.0,4.1,4.5,LGRAY,rounded=True)
_,tf=box(s,8.95,2.25,3.6,4.1)
P(tf,"VALIDACIÓN TEMPORAL",12.5,BLACK,bold=True,first=True,font=FH,sa=8)
P(tf,"Train = ediciones previas. Test = última encuesta 2025.2 (nunca vista).",12.5,BLACK,sa=10)
for t in ["Lineal simple: 0.13","Lineal múltiple: 0.22","Árbol: 0.20","Polinómica: 0.28"]:
    P(tf,t,12.5,GRAY,bullet=True,sa=6)
P(tf,"Random Forest: 0.27 (elegido)",13.5,BLACK,bold=True,bullet=True,sa=6)

# ============ 16 POR QUE RF
s=slide(BLACK); kicker(s,"05","Random Forest",True); title(s,"Por qué elegimos Random Forest",True)
why=[("A la par del mejor R²","empata con la polinómica en test (R²≈0.27); se elige por lo demás."),
     ("Capta no linealidad e interacciones","sin ingeniería manual (p. ej. el plateau de la experiencia)."),
     ("Robusto","tolera outliers y mezcla de escalas; no necesita estandarizar."),
     ("Sin overfitting","regularizado (profundidad y hojas mínimas): train ≈ test."),
     ("Interpretable","entrega importancia de variables para la lectura de negocio.")]
y=2.0
for i,(t,d) in enumerate(why):
    rect(s,0.7,y,0.55,0.55,ACC)
    _,tf=box(s,0.72,y+0.05,0.55,0.5); P(tf,"✓",18,BLACK,bold=True,align=PP_ALIGN.CENTER,first=True,sa=0)
    _,tf=box(s,1.45,y-0.03,11.2,0.95); P(tf,t,15,WHITE,bold=True,first=True,font=FH,sa=2); P(tf,d,13,GRAY)
    y+=0.98

# ============ 17 KMEANS
s=slide(WHITE); kicker(s,"05","K-Means",False); title(s,"Segmentación no supervisada",False)
img_fit(s,PRE/"clusters_bn.png",0.5,1.9,7.7,4.4)
rect(s,8.5,2.0,4.3,4.5,BLACK,rounded=True)
_,tf=box(s,8.75,2.25,3.8,4.1)
P(tf,"4 ARQUETIPOS",14,ACC,bold=True,first=True,font=FH,sa=10)
for t in ["Junior remoto (49%) — $2,75M","Semi-senior dolarizado (21%) — $2,81M",
          "Senior remoto (21%) — $4,43M · el mejor pago","Senior +15 corporativo estable (9%) — $4,01M"]:
    P(tf,t,13,WHITE,bullet=True,sa=10)
P(tf,"k=4 elegido con codo + silueta.",12.5,ACC,italic=True)
_,tf=box(s,0.5,6.45,7.9,0.9)
P(tf,"Nota: K-Means asigna cada punto a su centroide más cercano (verificado al 100%). Se ven "
     "solapados por la proyección 2D (muestra el 62% de la info) y porque los perfiles son un "
     "continuo — no por error del modelo.",10,GRAY,first=True,italic=True)

# ============ 18 DIFICULTADES
s=slide(WHITE); kicker(s,"05","Dificultades con el dataset",False); title(s,"Problemas y cómo se resolvieron",False)
dif=[("Comparar a través de 10× de inflación","deflación por IPC y, en USD, por inflación de EE.UU."),
     ("Ajuste según la moneda nativa","un sueldo dolarizado se ajusta por inflación de su país, no la local."),
     ("Seniority inferido en ediciones viejas","generaba multicolinealidad → se excluyó sin perder R²."),
     ("Validación realista","split temporal (no aleatorio) para simular el uso real."),
     ("Ruido de las tecnologías","sacadas del modelo (su premium era proxy del contexto); roles agrupados en top-15.")]
y=2.0
for t,d in dif:
    rect(s,0.7,y,0.5,0.5,BLACK)
    _,tf=box(s,0.7,y+0.04,0.5,0.45); P(tf,"!",18,ACC,bold=True,align=PP_ALIGN.CENTER,first=True,font=FH,sa=0)
    _,tf=box(s,1.4,y-0.03,11.3,0.95); P(tf,t,14.5,BLACK,bold=True,first=True,font=FH,sa=2); P(tf,d,12.5,GRAY)
    y+=0.98

# ============ 19 RESULTADOS
s=slide(WHITE); kicker(s,"06","Resultados y storytelling",False); title(s,"Brechas salariales y región",False)
img_fit(s,EDA/"eda_05_geografia.png",0.5,1.9,7.4,4.4)
rect(s,8.1,1.95,4.7,4.6,BLACK,rounded=True)
_,tf=box(s,8.4,2.2,4.1,4.2)
P(tf,"3×",44,ACC,bold=True,first=True,font=FH,sa=2)
P(tf,"brecha junior → senior: la experiencia manda",13,WHITE,sa=14)
P(tf,"Brechas salariales:",13.5,ACC,bold=True,font=FH,sa=4)
P(tf,"liderazgo +0,19 (la mayor)  ·  remoto +43% en USD",12.5,WHITE,sa=8)
P(tf,"Comparativa regional:",13.5,ACC,bold=True,font=FH,sa=4)
P(tf,"CABA +24%  ·  Patagonia (energía) destaca",12.5,WHITE,sa=8)
P(tf,"Dolarización (H5): no dio más estabilidad.",12.5,WHITE,italic=True)

# ============ 20 ARGENTINA VS MUNDO
s=slide(WHITE); kicker(s,"06","Resultados · comparativa internacional",False)
title(s,"Argentina vs el mundo",False)
img_fit(s,PRE/"comparacion_mundo.png",0.4,1.85,8.7,4.9)
rect(s,9.2,2.0,3.6,4.5,BLACK,rounded=True)
_,tf=box(s,9.45,2.25,3.15,4.1)
P(tf,"STACK OVERFLOW",13,ACC,bold=True,first=True,font=FH,sa=8)
for t in ["Argentina USD 3.333/mes","Arriba de Brasil, Ucrania e India","Abajo de Europa y EE.UU. (12.500)"]:
    P(tf,t,12.5,WHITE,bullet=True,sa=9)
P(tf,"Sysarmy mide USD 2.264: más conservador. Stack Overflow capta perfiles más "
     "internacionales / senior.",12,ACC,italic=True)

# ============ 22 CONCLUSIONES
s=slide(BLACK); kicker(s,"07","Conclusiones",True); title(s,"Hallazgos y valor estratégico",True)
_,tf=box(s,0.6,1.95,6.5,4.9)
P(tf,"Principales hallazgos",15,ACC,bold=True,first=True,font=FH,sa=8)
for t in ["El perfil explica ~27% del sueldo: experiencia y liderazgo mandan.",
          "Dolarizarse no dio más premium ni estabilidad.",
          "CABA, grandes empresas y backend (Go/Rust) concentran lo mejor.",
          "Pipeline 100% reproducible: fuentes crudas versionadas en parquet."]:
    P(tf,t,13.5,WHITE,bullet=True,sa=10)
rect(s,7.4,1.95,5.3,4.9,DCARD,rounded=True)
_,tf=box(s,7.7,2.25,4.7,4.4)
P(tf,"Oportunidades y valor",15,ACC,bold=True,first=True,font=FH,sa=8)
for t in ["Benchmarking salarial objetivo y por segmento.",
          "Lectura del costo del talento frente a la macro.",
          "Decisiones de oferta y retención basadas en datos.",
          "Herramienta viva: se reentrena con cada edición."]:
    P(tf,t,13.5,WHITE,bullet=True,sa=11)

prs.save(str(OUT)); print("PPTX:", OUT, "|", len(prs.slides._sldIdLst), "slides")
